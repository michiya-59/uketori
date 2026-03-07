"""
ウケトリ システム概要 PowerPointスライド生成スクリプト
python-pptx を使用して15枚のスライドを生成する
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# カラーパレット
BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1F, 0x29, 0x37)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
MEDIUM_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
GREEN = RGBColor(0x05, 0x96, 0x69)


def set_slide_bg(slide, color):
    """スライドの背景色を設定する"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    """背景用の矩形シェイプを追加する"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_text(slide, text, left, top, width, height, font_size=36, color=BLUE, bold=True, alignment=PP_ALIGN.LEFT):
    """タイトルテキストを追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_body_text(slide, text, left, top, width, height, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    """本文テキストを追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=18, color=DARK):
    """箇条書きリストを追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_table(slide, rows, cols, data, left, top, width, height, header_color=BLUE):
    """テーブルを追加する"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                if row_idx == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table_shape


def add_accent_line(slide, left, top, width):
    """アクセント線を追加する"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape


def create_presentation():
    """メインのプレゼンテーション生成関数"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== スライド1: 表紙 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # 左側の青い帯
    add_shape_bg(slide, Inches(0), Inches(0), Inches(0.4), Inches(7.5), BLUE)

    # メインタイトル
    add_title_text(slide, "ウケトリ（UKETORI）", Inches(1.5), Inches(1.5), Inches(10), Inches(1),
                   font_size=48, color=BLUE, bold=True)

    # タグライン
    add_title_text(slide, "見積から入金まで、ぜんぶウケトリ。", Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
                   font_size=32, color=DARK, bold=False)

    add_accent_line(slide, Inches(1.5), Inches(3.8), Inches(3))

    # サブタイトル
    add_body_text(slide, "AI搭載 受発注・請求・入金回収管理SaaS", Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
                  font_size=24, color=DARK)
    add_body_text(slide, "中小企業・フリーランス向け", Inches(1.5), Inches(4.9), Inches(10), Inches(0.6),
                  font_size=20, color=MEDIUM_GRAY)

    # ===== スライド2: 課題提起 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "中小企業の請求管理には「回収」の壁がある", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    problems = [
        "請求書を送っても入金確認が手作業で漏れが発生",
        "督促タイミングを逃し、回収遅延が常態化",
        "他ツールからの移行が面倒で乗り換えに踏み切れない",
        "Excelや紙管理ではリアルタイムな経営判断ができない",
    ]

    for i, problem in enumerate(problems):
        y = Inches(1.8) + Inches(i * 1.1)
        # アイコン代わりの丸
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.4), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RED
        shape.line.fill.background()
        add_body_text(slide, f"✕", Inches(1.05), y - Pt(2), Inches(0.4), Inches(0.4),
                      font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_body_text(slide, problem, Inches(1.7), y + Pt(2), Inches(9), Inches(0.4),
                      font_size=20, color=DARK)

    # 引用
    add_shape_bg(slide, Inches(0.8), Inches(6), Inches(11), Inches(0.8), LIGHT_BLUE)
    add_body_text(slide, "日本の中小企業の約60%が売掛金回収に課題を抱えている",
                  Inches(1.2), Inches(6.15), Inches(10), Inches(0.5), font_size=18, color=BLUE, bold=True)

    # ===== スライド3: ソリューション概要 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "ウケトリは「確実に回収する」ことにフォーカスしたSaaS",
                   Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8), font_size=30, color=WHITE)

    # フロー
    add_body_text(slide, "見積 → 受発注 → 納品 → 請求 → 入金回収 を一気通貫で管理",
                  Inches(0.8), Inches(1.6), Inches(11), Inches(0.5), font_size=22, color=DARK)

    add_title_text(slide, "2つの差別化軸", Inches(0.8), Inches(2.4), Inches(5), Inches(0.6),
                   font_size=26, color=BLUE)

    # 差別化カード1
    add_shape_bg(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(3.5), LIGHT_BLUE)
    add_title_text(slide, "① 入金回収特化", Inches(1.2), Inches(3.5), Inches(5), Inches(0.5),
                   font_size=24, color=BLUE)
    add_bullet_list(slide, [
        "• AI入金消込（5段階マッチング）",
        "• 自動督促メール送信",
        "• 回収率ダッシュボード",
        "• 与信スコアリング（0-100）",
    ], Inches(1.2), Inches(4.2), Inches(4.8), Inches(2.5), font_size=18)

    # 差別化カード2
    add_shape_bg(slide, Inches(7), Inches(3.2), Inches(5.5), Inches(3.5), LIGHT_BLUE)
    add_title_text(slide, "② 移行爆速", Inches(7.4), Inches(3.5), Inches(5), Inches(0.5),
                   font_size=24, color=BLUE)
    add_bullet_list(slide, [
        "• AI自動カラムマッピング",
        "• 5ステップ簡単ウィザード",
        "• 最短5分で移行完了",
        "• freee, Misoca, Excel等対応",
    ], Inches(7.4), Inches(4.2), Inches(4.8), Inches(2.5), font_size=18)

    # ===== スライド4: ターゲットユーザー =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "ターゲットユーザー", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    targets = [
        ("従業員 1〜50名", "中小企業・個人事業主"),
        ("業種不問", "全業種対応・業種別テンプレート搭載"),
        ("現在の管理方法", "Excel・紙・スプレッドシート・他SaaSで管理中"),
    ]

    for i, (title, desc) in enumerate(targets):
        y = Inches(1.8) + Inches(i * 1.3)
        add_shape_bg(slide, Inches(0.8), y, Inches(0.15), Inches(0.8), BLUE)
        add_body_text(slide, title, Inches(1.3), y, Inches(5), Inches(0.4),
                      font_size=22, color=DARK, bold=True)
        add_body_text(slide, desc, Inches(1.3), y + Inches(0.45), Inches(8), Inches(0.4),
                      font_size=18, color=MEDIUM_GRAY)

    # 最重要ターゲット
    add_shape_bg(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.3), RGBColor(0xFE, 0xF2, 0xF2))
    add_title_text(slide, "最重要ターゲット", Inches(1.2), Inches(5.6), Inches(5), Inches(0.4),
                   font_size=22, color=RED)
    add_body_text(slide, "「入金回収に課題がある」事業者 — 請求書を送って終わり、入金確認は月末に手作業、督促は気まずくて後回し…",
                  Inches(1.2), Inches(6.15), Inches(10.5), Inches(0.5), font_size=18, color=DARK)

    # ===== スライド5: 主要機能一覧 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "主要機能一覧", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    table_data = [
        ["カテゴリ", "機能"],
        ["帳票管理", "見積書・発注書・注文請書・納品書・請求書・領収書の作成〜PDF生成・メール送信"],
        ["案件管理", "商談〜入金完了までのパイプライン管理、ステータス自動遷移"],
        ["顧客マスタ", "顧客・取引先情報、担当者管理、タグ分類"],
        ["★ AI入金消込", "銀行CSV取込 → 5段階AIマッチングで自動消込"],
        ["★ 自動督促", "ルールベースの段階的督促メール自動送信"],
        ["★ 与信スコア", "支払い実績から取引先を0-100でスコアリング"],
        ["★ 回収ダッシュボード", "KPI・売掛金年齢表・入金予測"],
        ["★ データ移行", "AIカラムマッピングで他ツールから爆速移行"],
    ]

    add_table(slide, 9, 2, table_data, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))

    # ===== スライド6: AI入金消込 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "★ 差別化機能①：AI入金消込", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    add_title_text(slide, "5段階のインテリジェントマッチング", Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
                   font_size=24, color=BLUE)

    steps = [
        ("1", "ルールベースマッチ", "金額・日付の完全一致"),
        ("2", "名前正規化マッチ", "会社名の表記揺れを吸収（カ）→ 株式会社 等）"),
        ("3", "Claude AIマッチ", "AIが文脈から振込名義と請求書を推定"),
        ("4", "分類マッチ", "過去の消込パターンから学習"),
        ("5", "自動消込実行", "高確信度の結果を自動で消込処理"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.3) + Inches(i * 0.9)
        # 番号円
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        add_body_text(slide, num, Inches(1), y + Pt(2), Inches(0.5), Inches(0.5),
                      font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_body_text(slide, title, Inches(1.8), y + Pt(2), Inches(3), Inches(0.4),
                      font_size=20, color=DARK, bold=True)
        add_body_text(slide, desc, Inches(5), y + Pt(2), Inches(7), Inches(0.4),
                      font_size=18, color=MEDIUM_GRAY)

    # 引用
    add_shape_bg(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7), LIGHT_BLUE)
    add_body_text(slide, "銀行CSVをアップロードするだけで、手作業ゼロの入金消込を実現",
                  Inches(1.2), Inches(6.4), Inches(10.5), Inches(0.5), font_size=18, color=BLUE, bold=True)

    # ===== スライド7: 自動督促システム =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "★ 差別化機能②：自動督促システム", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    # 左側: 段階的督促
    add_title_text(slide, "段階的な督促ルール設定", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                   font_size=24, color=BLUE)

    dunning_steps = [
        ("第1段階（+3日）", "やんわりリマインドメール"),
        ("第2段階（+7日）", "正式な督促メール"),
        ("第3段階（+14日）", "強めの督促 + 上長CC"),
    ]

    for i, (stage, action) in enumerate(dunning_steps):
        y = Inches(2.3) + Inches(i * 1.0)
        add_shape_bg(slide, Inches(0.8), y, Inches(5.5), Inches(0.8), LIGHT_GRAY)
        add_body_text(slide, stage, Inches(1.2), y + Pt(4), Inches(3), Inches(0.3),
                      font_size=18, color=DARK, bold=True)
        add_body_text(slide, action, Inches(1.2), y + Pt(28), Inches(4.5), Inches(0.3),
                      font_size=16, color=MEDIUM_GRAY)

    # 右側: 与信スコア
    add_title_text(slide, "与信スコアリング（0-100点）", Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                   font_size=24, color=BLUE)

    score_data = [
        ["スコア帯", "評価", "対応"],
        ["80-100", "優良", "通常取引"],
        ["50-79", "注意", "取引条件見直し検討"],
        ["0-49", "警戒", "前払い要請・取引制限"],
    ]

    add_table(slide, 4, 3, score_data, Inches(7), Inches(2.3), Inches(5.5), Inches(2.5))

    # ===== スライド8: データ移行ウィザード =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "★ 差別化機能③：データ移行ウィザード", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    add_title_text(slide, "AI自動カラムマッピングで最短5分移行", Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
                   font_size=24, color=BLUE)

    wizard_steps = [
        ("1", "ソース選択", "移行元ツールを選択"),
        ("2", "ファイルアップロード", "CSV/ExcelをD&D"),
        ("3", "AIマッピング確認", "3段階AI自動対応付け"),
        ("4", "プレビュー", "変換後データ確認・修正"),
        ("5", "実行", "一括インポート+重複検出"),
    ]

    # 横並びのステップカード
    for i, (num, title, desc) in enumerate(wizard_steps):
        x = Inches(0.8) + Inches(i * 2.5)
        y = Inches(2.5)

        # カード背景
        add_shape_bg(slide, x, y, Inches(2.2), Inches(2.5), LIGHT_BLUE)

        # 番号
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y + Inches(0.2), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        add_body_text(slide, num, x + Inches(0.8), y + Inches(0.25), Inches(0.6), Inches(0.5),
                      font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # タイトル・説明
        add_body_text(slide, title, x + Inches(0.1), y + Inches(1.0), Inches(2), Inches(0.4),
                      font_size=18, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_body_text(slide, desc, x + Inches(0.1), y + Inches(1.5), Inches(2), Inches(0.8),
                      font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # 矢印
    for i in range(4):
        x = Inches(3.1) + Inches(i * 2.5)
        add_body_text(slide, "→", x, Inches(3.5), Inches(0.5), Inches(0.5),
                      font_size=28, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # 引用
    add_shape_bg(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.7), LIGHT_BLUE)
    add_body_text(slide, "他ツールからの乗り換え障壁をゼロにする",
                  Inches(1.2), Inches(5.9), Inches(10.5), Inches(0.5), font_size=18, color=BLUE, bold=True)

    # ===== スライド9: 帳票フロー全体像 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "帳票フロー全体像", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    # 帳票フローの各ステップ
    doc_types = [
        ("見積書", "estimate"),
        ("発注書", "PO"),
        ("納品書", "delivery"),
        ("請求書", "invoice"),
        ("領収書", "receipt"),
    ]

    for i, (jp, en) in enumerate(doc_types):
        x = Inches(0.5) + Inches(i * 2.5)
        y = Inches(2.0)

        color = RED if jp == "請求書" else BLUE
        add_shape_bg(slide, x, y, Inches(2.0), Inches(1.2), color)
        add_body_text(slide, jp, x + Inches(0.1), y + Inches(0.1), Inches(1.8), Inches(0.5),
                      font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_body_text(slide, en, x + Inches(0.1), y + Inches(0.6), Inches(1.8), Inches(0.4),
                      font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)

    # 矢印
    for i in range(4):
        x = Inches(2.6) + Inches(i * 2.5)
        add_body_text(slide, "→", x, Inches(2.3), Inches(0.5), Inches(0.5),
                      font_size=28, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # 入金消込ボックス
    add_shape_bg(slide, Inches(3.5), Inches(4.0), Inches(3), Inches(1.5), GREEN)
    add_body_text(slide, "AI入金消込", Inches(3.5), Inches(4.2), Inches(3), Inches(0.4),
                  font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_body_text(slide, "自動マッチング / 手動消込", Inches(3.5), Inches(4.7), Inches(3), Inches(0.4),
                  font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 自動督促ボックス
    add_shape_bg(slide, Inches(7.5), Inches(4.0), Inches(3), Inches(1.5), RED)
    add_body_text(slide, "自動督促", Inches(7.5), Inches(4.2), Inches(3), Inches(0.4),
                  font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_body_text(slide, "段階的メール送信", Inches(7.5), Inches(4.7), Inches(3), Inches(0.4),
                  font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 矢印テキスト
    add_body_text(slide, "↓ 請求後", Inches(4), Inches(3.4), Inches(2), Inches(0.4),
                  font_size=16, color=DARK, alignment=PP_ALIGN.CENTER)
    add_body_text(slide, "→ 未回収時", Inches(6.5), Inches(4.5), Inches(1.5), Inches(0.4),
                  font_size=16, color=DARK, alignment=PP_ALIGN.CENTER)

    # ワンクリック変換の説明
    add_shape_bg(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.7), LIGHT_GRAY)
    add_body_text(slide, "帳票間はワンクリックで変換可能（見積→請求、発注→納品 等）",
                  Inches(1.2), Inches(6.3), Inches(10.5), Inches(0.5), font_size=18, color=DARK)

    # ===== スライド10: システムアーキテクチャ =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "システムアーキテクチャ", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    arch_data = [
        ["レイヤー", "技術選定", "選定理由"],
        ["フロントエンド", "Next.js 16 + TypeScript + Tailwind + shadcn/ui", "App Router SSR/SSG、型安全、高品質UI"],
        ["バックエンドAPI", "Ruby on Rails 8.0 API mode", "高速開発、豊富なエコシステム"],
        ["データベース", "PostgreSQL 16", "JSONB、全文検索、堅牢性"],
        ["ジョブキュー", "SolidQueue（PostgreSQL内蔵）", "Redis不要でコスト削減"],
        ["キャッシュ", "SolidCache（PostgreSQL内蔵）", "Redis不要でコスト削減"],
        ["ファイル保存", "Cloudflare R2 / MinIO（S3互換）", "エグレス無料、低コスト"],
        ["認証", "JWT（自前実装）+ Pundit", "シンプル・軽量"],
        ["AI", "Claude API（Anthropic）", "消込・マッピング・帳票提案"],
    ]

    add_table(slide, 9, 3, arch_data, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))

    # ===== スライド11: インフラ構成 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "インフラ構成 — 極限コスト運用", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    infra_data = [
        ["サービス", "用途", "月額"],
        ["Vercel (Hobby→Pro)", "Next.js ホスティング", "¥0〜$20"],
        ["Fly.io (shared-cpu)", "Rails API + SolidQueue", "$3〜5"],
        ["Neon PostgreSQL", "データベース", "¥0"],
        ["Cloudflare R2", "PDF・画像保存", "¥0"],
        ["Cloudflare", "DNS + CDN + SSL", "¥0"],
        ["SendGrid", "メール送信", "¥0"],
        ["GitHub Actions", "CI/CD", "¥0"],
    ]

    add_table(slide, 8, 3, infra_data, Inches(0.8), Inches(1.5), Inches(8), Inches(4.5))

    # 合計コスト強調
    add_shape_bg(slide, Inches(9.5), Inches(2.5), Inches(3.3), Inches(2.5), LIGHT_BLUE)
    add_title_text(slide, "月額合計", Inches(9.5), Inches(2.7), Inches(3.3), Inches(0.5),
                   font_size=20, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_title_text(slide, "¥500〜3,000", Inches(9.5), Inches(3.3), Inches(3.3), Inches(0.8),
                   font_size=36, color=RED, alignment=PP_ALIGN.CENTER)
    add_body_text(slide, "Phase 1: 0〜100ユーザー", Inches(9.5), Inches(4.2), Inches(3.3), Inches(0.4),
                  font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # ===== スライド12: スケールアップ戦略 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "スケールアップ戦略", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    scale_data = [
        ["ユーザー数", "構成変更", "月額目安"],
        ["0〜50", "極限構成（現状）", "¥500〜3,000"],
        ["50〜200", "Neon Pro + Vercel Pro", "¥6,000〜10,000"],
        ["200〜500", "Fly.io 1GB RAM + Fly.io Postgres", "¥10,000〜15,000"],
        ["500〜1,000", "Fly.io 2台構成 + SolidQueue分離", "¥15,000〜25,000"],
        ["1,000+", "AWS移行（ECS + RDS）", "¥30,000〜"],
    ]

    add_table(slide, 6, 3, scale_data, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5))

    # 引用
    add_shape_bg(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.7), LIGHT_BLUE)
    add_body_text(slide, "PostgreSQL + S3互換APIでポータビリティを確保し、段階的にスケールアップ",
                  Inches(1.2), Inches(5.6), Inches(10.5), Inches(0.5), font_size=18, color=BLUE, bold=True)

    # ===== スライド13: 料金プラン =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "料金プラン", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    plan_data = [
        ["", "Free", "Starter", "Standard", "Professional"],
        ["ユーザー数", "3名", "5名", "10名", "無制限"],
        ["AI消込", "✕", "○", "○", "○"],
        ["自動督促", "✕", "✕", "○", "○"],
        ["データ移行", "✕", "○", "○", "○"],
        ["与信スコア", "✕", "✕", "○", "○"],
        ["API連携", "✕", "✕", "✕", "○"],
    ]

    add_table(slide, 7, 5, plan_data, Inches(1.5), Inches(1.5), Inches(10), Inches(4.5))

    add_shape_bg(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7), LIGHT_GRAY)
    add_body_text(slide, "Freeプランでも基本的な帳票管理は全機能利用可能",
                  Inches(1.2), Inches(6.4), Inches(10.5), Inches(0.5), font_size=18, color=DARK, bold=True)

    # ===== スライド14: 開発実績 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), BLUE)
    add_title_text(slide, "開発実績 — MVP 12週間で100%完了", Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                   font_size=32, color=WHITE)

    dev_data = [
        ["Week", "実装内容", "テスト数"],
        ["1-2", "環境構築・DB設計・認証・テナント", "319"],
        ["3-4", "自社情報設定・業種テンプレート・顧客マスタ", "489"],
        ["5-6", "品目・見積書・PDF生成・メール送信", "547"],
        ["7-8", "請求書・帳票変換・入金管理", "564"],
        ["9-10", "★AI消込・★督促・回収ダッシュボード", "619"],
        ["11-12", "★データ移行・ダッシュボード・テスト・バグ修正", "693"],
    ]

    add_table(slide, 7, 3, dev_data, Inches(0.8), Inches(1.5), Inches(8), Inches(4))

    # 右側KPI
    kpis = [
        ("693", "RSpecテスト全パス"),
        ("20", "データベーステーブル"),
        ("12", "サービスクラス"),
        ("10", "バックグラウンドジョブ"),
    ]

    for i, (num, label) in enumerate(kpis):
        y = Inches(1.8) + Inches(i * 1.2)
        add_shape_bg(slide, Inches(9.5), y, Inches(3.3), Inches(1.0), LIGHT_BLUE)
        add_title_text(slide, num, Inches(9.5), y + Pt(2), Inches(3.3), Inches(0.5),
                       font_size=32, color=BLUE, alignment=PP_ALIGN.CENTER)
        add_body_text(slide, label, Inches(9.5), y + Inches(0.55), Inches(3.3), Inches(0.3),
                      font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # ===== スライド15: まとめ =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)

    add_title_text(slide, "まとめ", Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
                   font_size=40, color=WHITE)

    add_title_text(slide, "ウケトリは「請求して終わり」を「確実に回収」に変える",
                   Inches(0.8), Inches(1.8), Inches(11), Inches(0.8), font_size=28, color=WHITE, bold=False)

    add_accent_line(slide, Inches(0.8), Inches(2.8), Inches(3))

    summary_items = [
        "AI入金消込 で手作業ゼロの入金確認",
        "自動督促 で回収漏れを防止",
        "与信スコア で取引リスクを可視化",
        "爆速データ移行 で乗り換え障壁ゼロ",
        "月額¥500〜 の極限コスト運用",
    ]

    for i, item in enumerate(summary_items):
        y = Inches(3.2) + Inches(i * 0.65)
        add_body_text(slide, f"✓  {item}", Inches(1.2), y, Inches(10), Inches(0.5),
                      font_size=22, color=WHITE)

    # タグライン
    add_shape_bg(slide, Inches(2), Inches(6.0), Inches(9), Inches(1.0), RGBColor(0x1D, 0x4E, 0xD8))
    add_title_text(slide, "見積から入金まで、ぜんぶウケトリ。",
                   Inches(2), Inches(6.15), Inches(9), Inches(0.7),
                   font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 保存
    output_path = "/Users/e0195/重要/uketori/output/presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
